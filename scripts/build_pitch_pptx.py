"""AuraView 2차 발표용 PPTX (12쪽) 생성.

16:9 와이드스크린. hero 슬라이드 = SDXL+PIL 합성 이미지 그대로,
정보 슬라이드 = 한글 폰트 + 표 + KPI 박스 네이티브 PowerPoint 도형.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "AuraView_pitch_2026.pptx"
CAPS = ROOT / "docs" / "captures"

# 16:9 와이드스크린 (1920x1080)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 색상
BG = RGBColor(0x08, 0x0C, 0x14)
SURFACE = RGBColor(0x0D, 0x15, 0x20)
LINE = RGBColor(0x1A, 0x23, 0x30)
TEXT = RGBColor(0xE8, 0xF0, 0xF8)
MUTED = RGBColor(0x8F, 0xA0, 0xB0)
ACCENT = RGBColor(0x00, 0xC8, 0xFF)
SAFE = RGBColor(0x00, 0xE0, 0x9A)
WARN = RGBColor(0xFF, 0xB0, 0x20)
BAD = RGBColor(0xFF, 0x6B, 0x6B)

FONT_KR = "맑은 고딕"


def set_font(run, name=FONT_KR, size=18, bold=False, color=TEXT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # eastAsia 폰트 명시 (한글 안정 렌더)
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", name)


def fill_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size=size, bold=bold, color=color)
    return tb


def add_eyebrow(slide, x, y, text):
    add_text(slide, x, y, Inches(8), Inches(0.3), text, size=12, bold=True, color=MUTED)


def add_h2(slide, x, y, w, text, color=TEXT):
    return add_text(slide, x, y, w, Inches(1.0), text, size=40, bold=True, color=color)


def add_box(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1)):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    box.adjustments[0] = 0.05
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = line_w
    box.shadow.inherit = False
    return box


def add_table(slide, x, y, w, h, headers, rows, header_color=ACCENT, body_size=14, header_size=12):
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers), Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h))).table

    # 헤더
    for i, hd in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE
        cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.1)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = hd
        set_font(r, size=header_size, bold=True, color=header_color)

    # 본문
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.1)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            color = TEXT
            if ri == len(rows) and headers[0] in ("도입 비율", "구간"):  # 합계행 강조
                color = SAFE
            set_font(r, size=body_size, color=color, bold=(ri == len(rows) and len(rows) >= 3))


def add_kpi(slide, x, y, w, h, val, label, color):
    add_box(slide, x, y, w, h, fill=BG, line=color, line_w=Pt(1.5))
    add_text(slide, x + Emu(int(Inches(0.2))), y + Emu(int(Inches(0.2))), w - Emu(int(Inches(0.4))), Inches(0.9),
             val, size=36, bold=True, color=color)
    add_text(slide, x + Emu(int(Inches(0.2))), y + Emu(int(Inches(1.0))), w - Emu(int(Inches(0.4))), Inches(0.4),
             label, size=12, color=MUTED)


def add_chip(slide, x, y, w, h, text, color=SAFE):
    add_box(slide, x, y, w, h, fill=BG, line=color, line_w=Pt(1))
    add_text(slide, x + Emu(int(Inches(0.15))), y + Emu(int(Inches(0.07))), w - Emu(int(Inches(0.3))), h,
             text, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)


def slide_hero(prs, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, BG)
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), 0, 0, width=SLIDE_W, height=SLIDE_H)
    return slide


# ============== SLIDES ==============

def s01_cover(prs):
    return slide_hero(prs, CAPS / "slide_01_cover.png")


def s02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    add_eyebrow(slide, Inches(0.6), Inches(0.4), "// 한국 도로의 현실")
    add_h2(slide, Inches(0.6), Inches(0.7), Inches(12.0), "매일 7명이 죽고, 800명이 다친다", color=TEXT)

    add_text(slide, Inches(0.6), Inches(1.7), Inches(6), Inches(0.5),
             "TAAS 2024 통계", size=22, bold=True, color=BAD)
    add_table(slide, Inches(0.6), Inches(2.3), Inches(6.0), Inches(3.4),
              ["지표", "값"],
              [["연간 사망", "2,581 명"],
               ["연간 부상", "290,400 명"],
               ["도시 교차로 사고 비중", "46%"],
               ["시야 가림 사고 비중", "22%"],
               ["고령 65세↑ 비중", "46%"]])

    add_text(slide, Inches(7.0), Inches(1.7), Inches(6), Inches(0.5),
             "기존 ADAS 한계", size=22, bold=True, color=WARN)
    lines = [
        "• 단일 카메라 시점 → 가려진 영역 사전 인지 불가",
        "• 선행경고 0.5~1초 → 회피 성공률 25% 미만",
        "• Tesla FSD 등 글로벌 솔루션은 한국 8 시나리오",
        "    (트럭·이륜·신호·우천·우회전·스쿨존·자전거·야간) 미반영",
        "• 민식이법 · 우회전 보행자 일시정지",
        "    → 운전자 형사 책임 강화",
    ]
    y0 = Inches(2.3)
    for i, ln in enumerate(lines):
        add_text(slide, Inches(7.0), y0 + Inches(0.45 * i), Inches(6), Inches(0.5),
                 ln, size=14, color=TEXT)
    add_text(slide, Inches(7.0), Inches(6.0), Inches(6), Inches(0.3),
             "출처: TAAS 2024 · KOTI 사회적 비용 추정", size=10, color=MUTED)


def s03_truck(prs):
    return slide_hero(prs, CAPS / "slide_03_truck.png")


def s04_fusion(prs):
    return slide_hero(prs, CAPS / "slide_04_fusion.png")


def s05_scenarios(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    add_eyebrow(slide, Inches(0.6), Inches(0.4), "// 한국 도로 특화 8 시나리오 × 도로교통법 매핑")
    add_h2(slide, Inches(0.6), Inches(0.7), Inches(12.0), "법령·판례·정량 기여 전부 매핑")

    headers = ["시나리오", "법령", "판례", "AuraView 기여"]
    rows = [
        ["① 트럭 가림",         "도교법 27조",            "2019도11622",    "occlusion +0.55"],
        ["② 좌측 사각 이륜",    "도교법 19조의2",         "2019도14517",    "측면 sweep"],
        ["③ 신호 가림",         "도교법 5조",             "2020도11458",    "신호 API + V2V"],
        ["④ 우천 교차로",       "도교법 19조 + 시행규칙", "2017도9534",     "환경 +0.45"],
        ["⑤ 우회전 보행자",     "도교법 25조 4항 (2022)", "2022도10752",    "sweep zone"],
        ["⑥ 스쿨존 (민식이법)", "도교법 12조 + 민식이법", "헌재 2019헌마927", "DSZ +0.62"],
        ["⑦ 자전거",            "도교법 13조 + 자전거법", "2021도8395",     "자전거 GIS +0.40"],
        ["⑧ 야간",              "도교법 48조",            "2018도12521",    "V2V 헤드라이트"],
    ]
    add_table(slide, Inches(0.6), Inches(1.7), Inches(12.0), Inches(5.0),
              headers, rows, body_size=13, header_size=12)
    add_text(slide, Inches(0.6), Inches(6.9), Inches(12), Inches(0.3),
             "출처: 국가법령정보센터 · 대법원 종합법률정보 · 헌법재판소 판례검색",
             size=10, color=MUTED)


def s06_schoolzone(prs):
    return slide_hero(prs, CAPS / "slide_05_schoolzone.png")


def s07_v2v(prs):
    return slide_hero(prs, CAPS / "slide_06_v2v.png")


def s08_ai_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    add_eyebrow(slide, Inches(0.6), Inches(0.4), "// AI 활용 — 학습도구 + 분석도구")
    add_h2(slide, Inches(0.6), Inches(0.7), Inches(12.0),
           "자체 학습 Risk Transformer · AUC 0.9403", color=TEXT)

    add_text(slide, Inches(0.6), Inches(1.8), Inches(6), Inches(0.4),
             "모델 사양", size=20, bold=True, color=ACCENT)
    add_table(slide, Inches(0.6), Inches(2.3), Inches(6.0), Inches(3.4),
              ["항목", "값"],
              [["구조", "Transformer (2-layer, d=64)"],
               ["파라미터", "67,970 개"],
               ["모델 크기", "278 KB (단말 임베드)"],
               ["학습 샘플", "10,000 (15 epoch)"],
               ["Optimizer", "AdamW (lr 1e-4)"]])

    add_text(slide, Inches(7.0), Inches(1.8), Inches(6), Inches(0.4),
             "학습 성능 (Validation)", size=20, bold=True, color=SAFE)
    kpis = [
        ("0.9403", "AUC (ROC)", ACCENT),
        ("0.9412", "F1 @ 0.5", ACCENT),
        ("0.9441", "Precision", SAFE),
        ("1.04 ms", "CPU p99", WARN),
    ]
    for i, (val, lbl, color) in enumerate(kpis):
        col, row = i % 2, i // 2
        x = Inches(7.0) + Emu(int(Inches(3.0) * col))
        y = Inches(2.3) + Emu(int(Inches(1.5) * row))
        add_kpi(slide, x, y, Inches(2.8), Inches(1.4), val, lbl, color)

    # AI 활용 증빙 chips
    add_chip(slide, Inches(7.0), Inches(5.4), Inches(2.4), Inches(0.4),
             "✓ AI 학습도구", SAFE)
    add_chip(slide, Inches(9.6), Inches(5.4), Inches(3.0), Inches(0.4),
             "✓ AI 분석도구 (ML Kit)", SAFE)

    add_text(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.3),
             "모델 가중치 + 학습 메트릭 별첨 제출 (models/risk_transformer.pt 278KB)",
             size=11, color=MUTED)


def s09_tesla(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    add_eyebrow(slide, Inches(0.6), Inches(0.4), "// 글로벌 솔루션 대비 한국 특화 5종")
    add_h2(slide, Inches(0.6), Inches(0.7), Inches(12.0), "Tesla FSD 시점 vs AuraView")

    headers = ["항목", "Tesla FSD 등 글로벌", "AuraView K-Perception"]
    rows = [
        ["차량 간 협업",   "자기 차량 시점만 인지",  "V2V Cross-Vehicle (heading 130° 가중)"],
        ["정류장 prior",   "일반 보행자 분류",        "Bus-Aware (정차/주행 prior +0.55)"],
        ["마주오는 차로",  "단방향 모델",             "Bidirectional + VDS 비대칭 분석"],
        ["공공 신호",      "Vision-only 인식",        "도로교통공단 신호 API + ITS 결합"],
        ["정책 환원",      "내부 데이터 폐쇄",        "위험 교차로 Top-N + DSZ 결합"],
    ]
    add_table(slide, Inches(0.6), Inches(2.0), Inches(12.0), Inches(4.8),
              headers, rows, body_size=15, header_size=13)
    add_text(slide, Inches(0.6), Inches(6.9), Inches(12), Inches(0.3),
             "상세: tesla_vs_auraview.svg · auraview.allthatai.kr/gallery/?f=비교",
             size=10, color=MUTED)


def s10_compliance(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    add_eyebrow(slide, Inches(0.6), Inches(0.4), "// 컴플라이언스 — 가명결합 + DSZ 안심구역")
    add_h2(slide, Inches(0.6), Inches(0.7), Inches(12.0), "개인정보보호법 · 국토부 훈령 완전 준수")

    add_text(slide, Inches(0.6), Inches(1.8), Inches(6), Inches(0.4),
             "가명결합 5단계 (개보법 28조의2)", size=20, bold=True, color=ACCENT)
    steps = [
        ("①", "가명화", "HMAC-SHA256 (식별자 별도 저장)"),
        ("②", "결합", "TAAS × VDS × 신호 3-table join"),
        ("③", "익명성", "k≥5 자동 검증"),
        ("④", "집계", "100m × 100m 그리드 셀"),
        ("⑤", "검증", "시점·k 값·결과 해시 자동 기록"),
    ]
    y0 = Inches(2.3)
    for i, (no, title, desc) in enumerate(steps):
        add_text(slide, Inches(0.6), y0 + Inches(0.55 * i), Inches(0.4), Inches(0.4),
                 no, size=18, bold=True, color=ACCENT)
        add_text(slide, Inches(1.1), y0 + Inches(0.55 * i), Inches(1.5), Inches(0.4),
                 title, size=14, bold=True, color=TEXT)
        add_text(slide, Inches(2.6), y0 + Inches(0.55 * i + 0.05), Inches(4.0), Inches(0.4),
                 desc, size=12, color=MUTED)

    add_text(slide, Inches(7.0), Inches(1.8), Inches(6), Inches(0.4),
             "DSZ 안심구역 (국토부 훈령 1456호)", size=20, bold=True, color=SAFE)
    add_table(slide, Inches(7.0), Inches(2.3), Inches(5.8), Inches(3.0),
              ["단계", "절차"],
              [["반입", "SHA-256 사전 검증"],
               ["결합", "안심구역 내 추가 결합"],
               ["분석", "위험 점수 + 우선순위"],
               ["반출", "통계만 (식별자 X)"],
               ["감사", "5년 로그 보존"]],
              header_color=SAFE, body_size=13)

    add_chip(slide, Inches(7.0), Inches(5.6), Inches(2.5), Inches(0.4),
             "✓ PII 마스킹 자동", SAFE)
    add_chip(slide, Inches(9.7), Inches(5.6), Inches(3.0), Inches(0.4),
             "✓ 100m 그리드 양자화", SAFE)


def s11_impact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    add_eyebrow(slide, Inches(0.6), Inches(0.4),
                "// 정량 사회 임팩트 (TAAS 2024 baseline · lead 3.38s · 회피율 84.5%)")
    add_h2(slide, Inches(0.6), Inches(0.7), Inches(12.5),
           "전국 도입 시 연 421명 사망 예방 · 5조 6,000억원 절감")

    add_table(slide, Inches(0.6), Inches(1.9), Inches(12.0), Inches(2.7),
              ["도입 비율", "사고 예방/년", "사망 감소", "부상 감소", "사회비용 절감"],
              [["Pilot 5%",   "1,694 건",  "21 명",  "2,370 명",  "약 2,800 억원"],
               ["확산 25%",   "8,470 건",  "105 명", "11,852 명", "약 1조 4,000 억원"],
               ["전국 100%",  "33,880 건", "421 명", "47,408 명", "약 5조 6,000 억원"]],
              body_size=15, header_size=13)

    add_text(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(0.4),
             "서울 위험 교차로 Top-10 우선 도입", size=18, bold=True, color=SAFE)
    add_text(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.5),
             "강남역 · 잠실역 · 광화문 · 신촌 · 청량리 · 건대입구 · 사당 · 홍대입구 · 신림 · 서울역 광장",
             size=13, color=TEXT)
    add_text(slide, Inches(0.6), Inches(5.95), Inches(12), Inches(0.5),
             "→ 합계 예방 효과 약 84명/년  ·  강남역 1곳 도입만 해도 연 11.8명 예방",
             size=14, bold=True, color=SAFE)

    add_text(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.3),
             "산출: TAAS 연간 × 도시교차로(46%) × 시나리오(42%) × 회피율(84.5%) × 도입률  ·  "
             "단가: KOTI 사회적 비용 2024 — 사망 5.5억/명, 중상 8천만/명",
             size=9, color=MUTED)


def s12_verify(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    add_eyebrow(slide, Inches(0.6), Inches(0.4),
                "// 1-step 라이브 검증 — 평가자 직접 확인 가능")
    add_h2(slide, Inches(0.6), Inches(0.7), Inches(12.5), "모두 지금 작동하는 시스템")

    # 좌 — 검증 URL
    add_text(slide, Inches(0.6), Inches(1.9), Inches(6), Inches(0.4),
             "검증 URL", size=20, bold=True, color=ACCENT)
    add_text(slide, Inches(0.6), Inches(2.4), Inches(8), Inches(0.5),
             "curl https://auraview.allthatai.kr/impact/submission-ready",
             size=12, color=SAFE)
    add_text(slide, Inches(0.6), Inches(2.9), Inches(8), Inches(0.5),
             "→ ready=true · 9/9 게이트 PASS", size=14, bold=True, color=TEXT)

    # 검증 strip
    chips = [
        ("✓ sources=25", SAFE),
        ("✓ schema v11", SAFE),
        ("✓ AUC 0.9403", SAFE),
        ("✓ git_sha live", SAFE),
        ("✓ banned 0 leaks", SAFE),
    ]
    x0 = Inches(0.6)
    for i, (txt, color) in enumerate(chips):
        add_chip(slide, x0 + Emu(int(Inches(1.4) * i)), Inches(3.6),
                 Inches(1.3), Inches(0.4), txt, color)

    # 페이지 chips
    pages = [
        "/ui 통합 대시보드",
        "/story 30초 일반인용",
        "/scorecard 평가 5종",
        "/policy 정책",
        "/competition 검증 허브",
    ]
    for i, p in enumerate(pages):
        col, row = i % 3, i // 3
        x = Inches(0.6) + Emu(int(Inches(2.2) * col))
        y = Inches(4.3) + Emu(int(Inches(0.55) * row))
        add_chip(slide, x, y, Inches(2.0), Inches(0.4), p, ACCENT)

    # 우 — 자료 패키지
    add_text(slide, Inches(7.5), Inches(1.9), Inches(5), Inches(0.4),
             "자료 패키지", size=20, bold=True, color=SAFE)
    add_table(slide, Inches(7.5), Inches(2.4), Inches(5.5), Inches(3.0),
              ["항목", "URL/사양"],
              [["라이브 시스템", "auraview.allthatai.kr"],
               ["오픈소스 (MIT)", "github.com/leelang7/AuraView"],
               ["네이티브 APK", "Galaxy Z Fold 3 v12.171 (34.6MB)"],
               ["AI 모델 가중치", "risk_transformer.pt 278KB"],
               ["별첨 PDF", "35쪽 (라이브 + 시각 + 근거)"]],
              header_color=SAFE, body_size=11)

    # B2C/B2B/B2G CTA box
    add_box(slide, Inches(7.5), Inches(5.6), Inches(5.5), Inches(1.4),
            fill=ACCENT, line=None)
    add_text(slide, Inches(7.7), Inches(5.7), Inches(5.3), Inches(0.5),
             "B2C · B2B · B2G", size=24, bold=True, color=BG)
    add_text(slide, Inches(7.7), Inches(6.2), Inches(5.3), Inches(0.5),
             "3년차 매출 70억 · 마이데이터 확장", size=13, bold=True, color=SURFACE)

    add_text(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.3),
             "2026 국토 · 교통 데이터 활용 경진대회 · AuraView K-Perception · v12.171",
             size=10, color=MUTED)


def main():
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

    print(f"[PPTX] building 12 슬라이드 ...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 12 슬라이드 빌드
    s01_cover(prs)
    s02_problem(prs)
    s03_truck(prs)
    s04_fusion(prs)
    s05_scenarios(prs)
    s06_schoolzone(prs)
    s07_v2v(prs)
    s08_ai_model(prs)
    s09_tesla(prs)
    s10_compliance(prs)
    s11_impact(prs)
    s12_verify(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"[OK] {OUT}")
    print(f"    {OUT.stat().st_size / 1024:.1f} KB · {len(prs.slides)} 슬라이드 · 16:9 와이드스크린")


if __name__ == "__main__":
    main()
